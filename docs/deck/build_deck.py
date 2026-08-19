"""Generates docs/deck/chakravyuh-walkthrough.pptx -- the mandatory "solution
walkthrough" submission artifact (issues.md I13). Every number in this deck
is pulled from a real project doc (docs/model-choice.md, docs/closed-loop.md,
docs/attack-catalogue.md, docs/research/) rather than invented for the slide
-- see the SOURCE comment above each content block. Re-run after any of
those documents change:

    uv run python docs/deck/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).resolve().parent / "chakravyuh-walkthrough.pptx"

NAVY = (0x0B, 0x1F, 0x3A)
ACCENT = (0xC8, 0x10, 0x2E)  # Mastercard-adjacent red, not the trademarked asset itself
LIGHT = (0xF5, 0xF6, 0xF8)


def _set_title_style(title_shape, size=32, color=NAVY):
    title_shape.text_frame.paragraphs[0].font.size = Pt(size)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = _rgb(color)


def _rgb(t):
    from pptx.dml.color import RGBColor

    return RGBColor(*t)


def add_bullet_slide(prs, title, bullets, *, subtitle=None, notes=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title)

    body = slide.placeholders[1].text_frame
    body.clear()
    first = True
    if subtitle:
        p = body.paragraphs[0]
        p.text = subtitle
        p.font.italic = True
        p.font.size = Pt(16)
        p.font.color.rgb = _rgb((0x55, 0x55, 0x55))
        first = False

    for level, text in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - level * 2)
        p.font.color.rgb = _rgb((0x22, 0x22, 0x22))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_slide(prs, title, headers, rows, *, notes=None, col_widths=None):
    layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title, size=28)

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.4), Inches(1.5), Inches(9.2), Inches(5.3)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(NAVY)
        cell.text_frame.paragraphs[0].font.color.rgb = _rgb((0xFF, 0xFF, 0xFF))

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(LIGHT)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_chart_slide(prs, title, categories, series_name, values, *, notes=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title, size=28)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    x, y, cx, cy = Inches(0.6), Inches(1.6), Inches(8.8), Inches(5.0)
    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Title -----------------------------------------------------------
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Chakravyuh"
    title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(48)
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = _rgb(NAVY)
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        "AI Defense Lab for Payment Security\n"
        "Mastercard Innovation Challenge 2026 -- Global Fintech Fest, Mumbai\n"
        "[Team name / members]"
    )

    # SOURCE: docs/master-project-brief.md section 1
    add_bullet_slide(
        prs,
        "The challenge",
        [
            (0, "One closed-loop system that plays both attacker and defender against GenAI-enabled payment fraud"),
            (1, "Identify -- map emerging GenAI-powered fraud vectors -- scored on diversity"),
            (1, "Generate -- simulate those attacks at scale as synthetic data -- scored on fidelity"),
            (1, "Defend -- detect them with an ML model -- scored on precision, recall, F1/AUC, low FPR"),
            (0, "Plus: novelty of the overall solution, real-world feasibility in live payments"),
            (0, "The loop is the point -- most submissions do “fake fraud → train model → 0.98 AUC,” which is circular and judges know it"),
        ],
        notes="Source: docs/master-project-brief.md section 1.",
    )

    # SOURCE: docs/master-project-brief.md section 2
    add_bullet_slide(
        prs,
        "Method: not brainstormed, derived",
        [
            (0, "1. Decompose each rail into numbered lifecycle steps: initiate → authenticate → authorise → settle → dispute"),
            (0, "2. Build the trust map -- who trusts whom, on what assumption, verified by what mechanism"),
            (0, "3. Break each trust anchor -- can the human be manipulated? Can the check be fooled, flooded, probed?"),
            (0, "4. Apply the actor lens -- outsider, insider, the customer themselves, an AI agent, a coordinated network"),
            (0, "5. Invert -- derive attacks that keep every conventional signal looking normal. This is where novelty lives."),
            (0, "Verification-strength ladder: V0 (nothing verifies it) and V1 (checked once, never re-run) are the attack surface"),
        ],
        notes="Source: docs/master-project-brief.md section 2.",
    )

    # SOURCE: docs/master-project-brief.md section 3.1
    add_bullet_slide(
        prs,
        "Three findings that shape the whole submission",
        [
            (0, "The UPI PIN proves the wrong proposition"),
            (1, "Cryptographically sound -- proves the right person pressed the keys. Says nothing about whether they understood or freely chose to. A V3 mechanism sitting on a V0 inference."),
            (0, "UPI credits are irreversible -- detection must be pre-authorisation"),
            (1, "No chargeback. A post-hoc detector is a reporting tool, not a control."),
            (0, "Authorised-but-deceived falls outside zero liability"),
            (1, "RBI's framework covers unauthorised transactions. A push payment made under deception is, on the record, authorised."),
        ],
        notes="Source: docs/master-project-brief.md section 3.1.",
    )

    # SOURCE: docs/attack-catalogue.md generator merge map
    add_table_slide(
        prs,
        "Diversity: 58 catalogue entries → 13 generators",
        ["ID", "Generator", "Data signature"],
        [
            ("G01", "scam_induced_push", "Genuine device/PIN, new beneficiary, coercion session fields"),
            ("G02", "mule_network", "Fan-in, pass-through, throughput ratio → 1"),
            ("G03", "card_testing_probe", "Micro-amounts, high decline rate, card rotation"),
            ("G04", "adversarial_evasion", "Probe then optimise; features tuned inside legit distribution"),
            ("G05", "first_party_dispute", "Transaction genuine; signal is claimant dispute history"),
            ("G06", "stealth_mandate", "Uniform small recurring amounts, max_amount ≫ actual_amount"),
            ("G07", "synthetic_merchant", "New merchant, step volume curve, registry-unverified KYB"),
            ("G08", "transaction_laundering", "mcc_declared ≠ mcc_inferred_from_basket"),
            ("G09", "credential_takeover", "Auth succeeds, device/session anomalous"),
            ("G10", "synthetic_identity_bustout", "Credit-building then simultaneous utilisation spike"),
            ("G11", "subthreshold_fragmentation", "Uniform amounts just under a limit"),
            ("G12", "agentic_injection", "is_agent_initiated, beneficiary ≠ seller of record"),
            ("G13", "insider_abuse", "No external anomaly; approval-velocity signals only"),
        ],
        col_widths=[0.7, 2.3, 6.2],
        notes="Source: docs/attack-catalogue.md, 'Generator merge map -- 58 entries → 13 generators'. Two attacks producing the same data signature are one generator, not two.",
    )

    # SOURCE: docs/attack-catalogue.md inversion pass
    add_bullet_slide(
        prs,
        "Novelty: attacks that stay inside the legitimate distribution",
        [
            (0, "Conventional detection watches: velocity, new beneficiary, amount thresholds, geo, device change, time-of-day, IP reputation"),
            (0, "These attacks keep all of them normal:"),
            (1, "STEALTH mandate -- amounts below every threshold, beneficiary is a merchant not a new payee"),
            (1, "FIRSTPARTY dispute -- the transaction is genuine; nothing to detect at transaction time at all"),
            (1, "DIGITALARREST -- own device, own PIN, correct entry; only session-level coercion fields and mule graph carry signal"),
            (0, "In a scam-induced payment everything is genuine -- the only anomaly is intent, which appears in no standard column"),
        ],
        notes="Source: docs/attack-catalogue.md 'Inversion pass' section.",
    )

    # SOURCE: docs/data-schema-v1.md
    add_bullet_slide(
        prs,
        "The schema: only fields a live system would actually have",
        [
            (0, "8 tables: transactions, parties, devices, merchants, mandates, disputes, graph_edges, labels"),
            (0, "Design rule: if an issuer/PSP wouldn't hold it at the moment of scoring, it doesn't belong"),
            (0, "Highest-value fields -- the only place the coercion signal becomes measurable:"),
            (1, "time_on_confirm_screen_s, screen_share_active, call_active_during_txn, accessibility_service_active, beneficiary_added_ago_s"),
            (0, "Strongest mule discriminators:"),
            (1, "has_salary_credit, organic_spend_ratio, throughput_ratio_24h, plus fan-in/fan-out from graph_edges"),
        ],
        notes="Source: docs/data-schema-v1.md.",
    )

    # SOURCE: issues.md I8, data/reference/*.json
    add_bullet_slide(
        prs,
        "Fidelity: validated against published reference datasets",
        [
            (0, "Generated background traffic compared against IEEE-CIS, PaySim, BankSim, ULB reference marginals -- not fabricated benchmarks"),
            (0, "Realistic: hour-of-day shape, right-skewed amount distribution (CV 2.18), heavy-tailed graph degree"),
            (0, "Two gaps found and logged honestly, not hidden:"),
            (1, "Day-of-week distribution is currently flat -- no weekend/salary-day effect yet (issues.md I15)"),
            (1, "Amount medians don't yet vary by MCC -- grocery, fuel and hotels cluster near the same value (issues.md I16)"),
            (0, "This is what “fidelity of the background, not the attacks” means in practice -- see docs/master-project-brief.md section 8"),
        ],
        notes="Source: issues.md I8/I15/I16, generated by src/validation/report.py against data/reference/.",
    )

    # SOURCE: docs/model-choice.md, evaluation methodology
    add_bullet_slide(
        prs,
        "Detection: evaluation methodology",
        [
            (0, "Headline metric: precision/recall at fixed FPR (0.1%, 1%) -- not ROC-AUC"),
            (1, "UPI credits are final; a detector is a pre-auth control, the operating point matters more than ranking quality"),
            (0, "Temporal split only -- train/validate/test by transaction timestamp, never random"),
            (1, "A random split leaks campaign structure across the boundary and inflates every number"),
            (0, "One attack family (synthetic_identity_bustout) held out entirely from training"),
            (1, "Tests generalisation to a genuinely unseen attack, not memorisation of the 13 generators trained on"),
            (0, "Every generator also emits its legit_lookalike population -- without it, the classifier separates two trivially different distributions"),
        ],
        notes="Source: docs/master-project-brief.md section 6, docs/model-choice.md.",
    )

    # SOURCE: docs/model-choice.md, honesty story
    add_bullet_slide(
        prs,
        "The honesty story: 3 leaks found and fixed, not papered over",
        [
            (0, "First result: PR-AUC 0.999, 100% recall everywhere, including the held-out family -- suspicious, and we said so"),
            (0, "Investigated instead of shipping it. Found and fixed three compounding issues (issues.md I6/I7/I17):"),
            (1, "I6 -- legit_lookalike rows were shallow copies of the fraud row (same counterparty, same timestamp)"),
            (1, "I7 -- attack campaigns hammered one fixed new counterparty; campaign shape alone was a near-perfect tell"),
            (1, "I17 -- a hardcoded ip_asn default meant every attack row carried the same ASN"),
            (0, "“An honest 0.82 beats a suspicious 0.99” -- docs/master-project-brief.md section 8"),
        ],
        notes="Source: issues.md I6/I7/I17, docs/model-choice.md.",
    )

    # SOURCE: docs/model-choice.md, final numbers
    add_table_slide(
        prs,
        "Detection results, after the fix",
        ["Metric", "Value"],
        [
            ("PR-AUC", "0.9866"),
            ("ROC-AUC (secondary)", "0.9997"),
            ("Precision @ 0.1% FPR", "0.8406"),
            ("Recall @ 0.1% FPR", "0.9775"),
            ("Precision @ 1% FPR", "0.4384"),
            ("Recall @ 1% FPR", "0.9902"),
            ("F1-optimal alerts / 1,000 txns", "4.99"),
            ("Held-out family recall @ 0.1%/1% FPR", "1.0000 (see caveat, next slide)"),
        ],
        col_widths=[5.5, 3.7],
        notes="Source: docs/model-choice.md, stage5/models/model_metadata.json.",
    )

    add_chart_slide(
        prs,
        "Feature importance: spread, not one dominant proxy",
        ["edge_count", "beneficiary_added_ago_s", "edge_value_total", "is_two_hop_passthrough", "other (long tail)"],
        "Importance",
        (0.34, 0.15, 0.06, 0.046, 0.404),
        notes=(
            "Source: docs/model-choice.md. Graph-based mule discriminators dominate -- "
            "exactly what docs/master-project-brief.md section 6 names as the strongest signal, not an artifact. "
            "Contrast with the pre-fix run where a single feature held 92.5% importance."
        ),
    )

    # SOURCE: docs/closed-loop.md
    add_bullet_slide(
        prs,
        "Closing the loop: detector → new attack → detector",
        [
            (0, "Diagnosis: the fixed detector leans on graph relationships (edge_count, beneficiary age) above everything else"),
            (0, "New attack variant: adversarial_evasion's adaptive mode targets exactly those two features"),
            (1, "Routes every event through the payer's single busiest existing relationship instead of a small pool"),
            (1, "Pushes beneficiary age further from “recently added” toward the population's typical maximum"),
            (0, "The mechanism is live: stage5/training/build_adaptive_attack_config.py reads the currently-saved model's feature_importances_ and derives this config automatically for the next generation run"),
            (0, "Two families are structurally exempt, and we say so: first_party_dispute and insider_abuse have no transaction-time signal by design -- see docs/closed-loop.md"),
        ],
        notes="Source: docs/closed-loop.md.",
    )

    # SOURCE: docs/research/d3-regulatory-limits.md
    add_table_slide(
        prs,
        "Real-world feasibility: current, sourced regulatory limits",
        ["Limit", "Value", "Source"],
        [
            ("UPI Lite per-transaction", "₹1,000", "RBI RDP statement 6 Dec 2024; NPCI OC 169-A/FY2024-25"),
            ("UPI Lite wallet balance", "₹5,000", "same"),
            ("Min-KYC PPI monthly load / balance", "₹10,000 / ₹10,000", "RBI Master Direction on PPIs, 2021"),
            ("UPI AutoPay AFA threshold (general)", "₹15,000", "RBI E-Mandate Framework 2026 (21 Apr 2026)"),
            ("UPI AutoPay AFA exemption (3 categories)", "₹1,00,000", "same -- mutual funds, insurance, credit-card bills only"),
        ],
        col_widths=[3.2, 2.0, 4.0],
        notes=(
            "Source: docs/research/d3-regulatory-limits.md. UPI P2M higher-limit circular number is a "
            "secondary-sourced estimate, not primary-confirmed -- flagged, not hidden. Do not quote it as confirmed."
        ),
    )

    # SOURCE: docs/attack-catalogue.md confidence tiers, evidence pass
    add_bullet_slide(
        prs,
        "Confidence, labelled honestly",
        [
            (0, "Every catalogue entry is tagged observed / emerging / speculative -- not overclaimed"),
            (0, "Evidence pass (docs/research/evidence-pass.md) verified citations for the 15 highest-confidence entries against real, dated sources"),
            (1, "Some downgraded on review rather than left overstated -- see that file for specifics"),
            (0, "D2 defect: corrected a wrong headline statistic (₹1,750cr claimed → ₹120.30cr actual, I4C Jan–Apr 2024 digital-arrest losses) before it reached this deck"),
            (0, "“A judge respects a clearly-marked speculative attack far more than an overclaimed one” -- docs/master-project-brief.md section 8"),
        ],
        notes="Source: docs/attack-catalogue.md fix log and confidence tiers, docs/research/evidence-pass.md.",
    )

    # SOURCE: web/app.py
    add_bullet_slide(
        prs,
        "Working prototype: the loop, live",
        [
            (0, "Streamlit app (web/app.py) -- uv run streamlit run web/app.py"),
            (0, "Live scoring -- pick or hand-tune a transaction, score it through the real Stage 6 pipeline: fraud probability, predicted attack family, contributing signals, analyst narrative"),
            (0, "Closed loop -- current feature importances, and what the next adaptive attack generation will target, read live from the saved model"),
            (0, "Attack catalogue -- the 13 generators, browsable"),
            (0, "[Insert screenshot once a model is trained in the presenting environment]"),
        ],
        notes="Source: web/app.py, web/README.md.",
    )

    # SOURCE: issues.md Next tasks
    add_bullet_slide(
        prs,
        "What's left, tracked openly",
        [
            (0, "issues.md is the live source of truth -- defects found, fixed, and still open, not swept into slides"),
            (1, "I11: retrain-and-measure step for the closed loop (heavy, deferred, mechanism is built and tested)"),
            (1, "I15/I16: weekday and MCC-conditioned amount calibration in the legitimate generator"),
            (1, "I9: attack-family classifier artifacts need regenerating against the current temporal split"),
            (0, "This is a deliberate choice: an honest punch list is more credible than a submission that claims to be finished"),
        ],
        notes="Source: issues.md.",
    )

    # Closing
    add_bullet_slide(
        prs,
        "Thank you",
        [
            (0, "Repository: github.com/Sharkyii/Chakravyuh"),
            (0, "docs/master-project-brief.md, docs/attack-catalogue.md, docs/model-choice.md, docs/closed-loop.md, issues.md"),
            (0, "[Team contact]"),
        ],
    )

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
